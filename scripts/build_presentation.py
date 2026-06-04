#!/usr/bin/env python3
"""Generate 扭蛋樂園 final presentation from Slidesgo template."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "docs" / "Multimedia Software Pitch Deck by Slidesgo.pptx"
OUTPUT = ROOT / "docs" / "扭蛋樂園_期末專題簡報.pptx"

FONT = "Microsoft JhengHei"
CONTENT_LAYOUT = "TITLE_AND_BODY"
TWO_COL_LAYOUT = "TITLE_AND_TWO_COLUMNS"


def delete_all_slides(prs):
    xml = prs.slides._sldIdLst
    for sld_id in list(xml):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        xml.remove(sld_id)


def layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout not found: {name}")


def placeholder(slide, ph_type, idx=None):
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph = shape.placeholder_format
        if ph.type != ph_type:
            continue
        if idx is None or ph.idx == idx:
            return shape
    return None


def _prep_text_frame(tf, wrap=True, auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE):
    tf.word_wrap = wrap
    try:
        tf.auto_size = auto_size
    except Exception:
        pass


def _body_font_size(line_count, max_chars=0):
    """Pick a size that fits the template body box (~8.4\" × 3.7\")."""
    if line_count <= 3 and max_chars < 42:
        return Pt(22)
    if line_count <= 5 and max_chars < 50:
        return Pt(18)
    if line_count <= 7:
        return Pt(16)
    return Pt(14)


def set_text(shape, text, size=Pt(18), bold=False, align=None):
    if shape is None:
        return
    tf = shape.text_frame
    _prep_text_frame(tf)
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold


def set_multiline(shape, lines, size=None, bold=False, bullet=True):
    if shape is None or not lines:
        return
    if size is None:
        max_len = max(len(ln) for ln in lines)
        size = _body_font_size(len(lines), max_len)

    tf = shape.text_frame
    _prep_text_frame(tf)
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = 1.15
        if bullet and line and not line.startswith("　") and not line.startswith("→"):
            # Template-friendly bullet character
            text = f"• {line}" if not line.startswith("•") else line
        else:
            text = line
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = size
        run.font.bold = bold and i == 0


def add_body_slide(prs, title, lines, body_size=None):
    slide = prs.slides.add_slide(layout_by_name(prs, CONTENT_LAYOUT))
    set_text(placeholder(slide, PP_PLACEHOLDER.TITLE, 0), title, Pt(28), True)
    set_multiline(
        placeholder(slide, PP_PLACEHOLDER.BODY, 1),
        lines,
        size=body_size,
    )
    return slide


def add_section(prs, title, number):
    slide = prs.slides.add_slide(layout_by_name(prs, "SECTION_HEADER"))
    set_text(placeholder(slide, PP_PLACEHOLDER.TITLE, 0), title, Pt(32), True)
    set_text(placeholder(slide, PP_PLACEHOLDER.TITLE, 2), number, Pt(36), True)
    return slide


def add_cover(prs, title, subtitle_lines):
    slide = prs.slides.add_slide(layout_by_name(prs, "TITLE"))
    set_text(placeholder(slide, PP_PLACEHOLDER.CENTER_TITLE, 0), title, Pt(36), True)
    sub = placeholder(slide, PP_PLACEHOLDER.SUBTITLE, 1)
    if isinstance(subtitle_lines, str):
        subtitle_lines = subtitle_lines.split("\n")
    set_multiline(sub, subtitle_lines, size=Pt(18), bullet=False)


def add_toc(prs, items):
    slide = prs.slides.add_slide(layout_by_name(prs, "BLANK_1_1_1_1_1_1"))
    set_text(placeholder(slide, PP_PLACEHOLDER.TITLE, 0), "簡報大綱", Pt(26), True)
    number_by_idx = {}
    labels_by_idx = {}
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph = shape.placeholder_format
        if ph.type == PP_PLACEHOLDER.TITLE and ph.idx != 0:
            number_by_idx[ph.idx] = shape
        elif ph.type == PP_PLACEHOLDER.SUBTITLE:
            labels_by_idx[ph.idx] = shape
    num_idxs = sorted(number_by_idx.keys())
    lbl_idxs = sorted(labels_by_idx.keys())
    for i, (num, label) in enumerate(items[:6]):
        if i < len(num_idxs):
            set_text(number_by_idx[num_idxs[i]], num, Pt(22), True)
        if i < len(lbl_idxs):
            # Short labels to avoid overflow in small TOC cells
            set_text(labels_by_idx[lbl_idxs[i]], label, Pt(14))
    return slide


def add_two_columns(prs, slide_title, left_title, left_lines, right_title, right_lines):
    slide = prs.slides.add_slide(layout_by_name(prs, TWO_COL_LAYOUT))
    set_text(placeholder(slide, PP_PLACEHOLDER.TITLE, 0), slide_title, Pt(26), True)

    # Template placeholder idx mapping (verified against layout geometry)
    left_hdr = placeholder(slide, PP_PLACEHOLDER.SUBTITLE, 3)
    right_hdr = placeholder(slide, PP_PLACEHOLDER.SUBTITLE, 4)
    left_body = placeholder(slide, PP_PLACEHOLDER.SUBTITLE, 2)
    right_body = placeholder(slide, PP_PLACEHOLDER.SUBTITLE, 1)

    set_text(left_hdr, left_title, Pt(18), True)
    set_text(right_hdr, right_title, Pt(18), True)
    set_multiline(left_body, left_lines, size=Pt(15))
    set_multiline(right_body, right_lines, size=Pt(15))
    return slide


def add_code_slide(prs, title, code_lines, explain_lines):
    add_body_slide(prs, title, code_lines, body_size=Pt(14))
    if explain_lines:
        add_body_slide(
            prs,
            f"{title}（說明）",
            explain_lines,
            body_size=Pt(18),
        )


def _chunk(lines, n):
    for i in range(0, len(lines), n):
        yield lines[i : i + n]


def build():
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template missing: {TEMPLATE}")

    prs = Presentation(str(TEMPLATE))
    delete_all_slides(prs)

    add_cover(prs, "扭蛋樂園", ["資料庫 × 全端整合期末專題", "Node.js · Express · SQLite"])

    add_toc(
        prs,
        [
            ("01", "專案與評分"),
            ("02", "ER 圖"),
            ("03", "SQL 範例"),
            ("04", "真實資料"),
            ("05", "現場展示"),
            ("06", "口試準備"),
        ],
    )

    add_body_slide(
        prs,
        "專案簡介",
        [
            "仿 CS 開箱的扭蛋轉盤網站",
            "註冊、儲值、抽獎、背包、售回、出貨",
            "Express + SQLite + 純前端",
            "7 張資料表 + 管理後台",
        ],
    )

    add_body_slide(
        prs,
        "評分項目對應",
        [
            "前端 10%：8 頁、表單、轉盤",
            "資料庫 20%：7 表、FK、正規化",
            "整合 20%：API 寫入 SQLite",
            "口試 50%：ER、SQL、真實資料",
        ],
    )

    add_section(prs, "系統架構", "01")

    add_body_slide(
        prs,
        "全端資料流",
        [
            "瀏覽器 → fetch /api（Session）",
            "Express 路由 + 參數化 SQL",
            "SQLite gachapon.db",
            "口試：按鈕 → 哪張表多一列",
        ],
    )

    add_body_slide(
        prs,
        "功能一覽",
        [
            "登入（bcrypt + session）",
            "錢包與交易流水",
            "伺服器加權抽獎",
            "售回 60%、多選出貨",
            "物流四階段追蹤",
        ],
    )

    add_section(prs, "ER 實體關聯圖", "02")

    add_body_slide(
        prs,
        "七個實體",
        [
            "users、machines、toys",
            "user_inventory、transactions",
            "shipments、shipment_items",
            "詳圖：docs/ER-diagram.md",
        ],
    )

    add_two_columns(
        prs,
        "關聯與基數",
        "Relationship",
        [
            "使用者 1:N 背包",
            "使用者 1:N 交易",
            "扭蛋機 1:N 公仔",
            "出貨 M:N（明細表）",
        ],
        "設計理由",
        [
            "款式 vs 實例分表",
            "transactions 流水帳",
            "UNIQUE 防重複出貨",
            "db.transaction()",
        ],
    )

    add_section(prs, "SQL 查詢範例", "03")

    add_code_slide(
        prs,
        "Q2：我的背包",
        [
            "SELECT ui.inventory_id, t.name,",
            "       t.rarity, t.price,",
            "       m.name AS machine_name",
            "FROM user_inventory ui",
            "JOIN toys t ON t.toy_id = ui.toy_id",
            "JOIN machines m ON m.machine_id = t.machine_id",
            "WHERE ui.user_id = ?",
            "  AND ui.status = 'owned';",
        ],
        [
            "JOIN 避免冗餘欄位",
            "只顯示持有中公仔",
            "? 由 session 代入",
        ],
    )

    add_code_slide(
        prs,
        "Q2：出貨追蹤",
        [
            "SELECT s.*,",
            "  COUNT(si.shipment_item_id)",
            "    AS item_count",
            "FROM shipments s",
            "LEFT JOIN shipment_items si",
            "  ON si.shipment_id = s.shipment_id",
            "WHERE s.user_id = ?",
            "GROUP BY s.shipment_id;",
        ],
        [
            "LEFT JOIN 彙總件數",
            "明細另查四表 JOIN",
        ],
    )

    add_section(prs, "真實資料示範", "04")

    add_body_slide(
        prs,
        "Q3：三筆真實資料",
        [
            "① users：帳號與餘額",
            "② transactions：儲值紀錄",
            "③ inventory JOIN toys：抽到的公仔",
            "以 user_id / toy_id 串聯",
        ],
    )

    add_code_slide(
        prs,
        "終端機查詢",
        [
            "sqlite3 db/gachapon.db",
            "SELECT * FROM users",
            "  WHERE username='帳號';",
            "SELECT * FROM transactions",
            "  WHERE user_id=1;",
        ],
        [
            "演示前：註冊→儲值→抽蛋",
            "deposit 的 toy_id 為 NULL",
        ],
    )

    add_section(prs, "現場功能展示", "05")

    add_body_slide(
        prs,
        "演示順序",
        [
            "npm start → localhost:3000",
            "註冊、儲值",
            "抽扭蛋、售回",
            "出貨與追蹤",
            "（選）管理後台改狀態",
        ],
    )

    add_two_columns(
        prs,
        "教授可能追問",
        "抽獎",
        [
            "邏輯在伺服器",
            "weight 加權隨機",
            "transaction 原子性",
        ],
        "出貨",
        [
            "shipment_items",
            "解開 M:N",
            "超商 20 件加倍",
        ],
    )

    add_section(prs, "口試追問準備", "06")

    add_body_slide(
        prs,
        "正規化 · 安全 · SQLite",
        [
            "公仔名稱只在 toys",
            "bcrypt + 參數化查詢",
            "pull/sell/ship 用 transaction",
            "migrate.sql 遷移",
        ],
    )

    add_cover(prs, "謝謝聆聽", ["Q1 ER 圖 · Q2 SQL · Q3 真實資料", "歡迎提問"])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved {len(prs.slides)} slides → {OUTPUT}")


if __name__ == "__main__":
    build()
